"""Slide rendering services — PPTX and PDF.

Extracted from routes/slides.py to follow Single Responsibility Principle.
Each renderer is a stateless class with a single `generate(items, title) -> bytes` method.
"""

from __future__ import annotations

import io
import logging
import math
from typing import TYPE_CHECKING, List

logger = logging.getLogger(__name__)

# ── python-pptx ───────────────────────────────────────────────────────────────

PPTX_OK = False
PPTX_ERROR: str = ""

try:
    from pptx import Presentation as _Presentation          # type: ignore
    from pptx.util import Inches as _Inches, Pt as _Pt      # type: ignore
    from pptx.dml.color import RGBColor as _RGBColor        # type: ignore
    from pptx.oxml import parse_xml as _parse_xml           # type: ignore
    _test = _Presentation()
    _buf  = io.BytesIO()
    _test.save(_buf)
    if _buf.tell() > 0:
        PPTX_OK = True
    del _test, _buf
except Exception as _e:
    PPTX_ERROR = str(_e)
    logger.warning("python-pptx unavailable: %s", _e)


# ── reportlab ────────────────────────────────────────────────────────────────

PDF_OK = False
PDF_ERROR: str = ""

try:
    from reportlab.lib.pagesizes import landscape, A4  # type: ignore
    from reportlab.lib.colors import HexColor           # type: ignore
    from reportlab.pdfgen.canvas import Canvas as _Canvas  # type: ignore
    PDF_OK = True
except Exception as _pe:
    PDF_ERROR = str(_pe)
    logger.warning("reportlab unavailable: %s", _pe)


# ── SlideItem protocol (avoid circular import) ────────────────────────────────

if TYPE_CHECKING:
    from ..routes.slides import SlideItem


# ── Shared file-download helper ───────────────────────────────────────────────

def make_download_headers(title: str, extension: str, slide_count: int) -> dict:
    """Build RFC 5987-compliant Content-Disposition headers (no latin-1 issues)."""
    import re
    from urllib.parse import quote
    safe    = re.sub(r"[^a-zA-Z0-9\s_-]", "", title).strip().replace(" ", "_") or "bai_giang"
    encoded = quote(title.strip() or "bai_giang") + f".{extension}"
    return {
        "Content-Disposition": f'attachment; filename="{safe}.{extension}"; filename*=UTF-8\'\'{encoded}',
        "X-Slide-Count": str(slide_count),
    }


# ── PPTX Renderer ─────────────────────────────────────────────────────────────

class PptxRenderer:
    """Generates .pptx bytes from a list of SlideItem objects."""

    C_BG     = (255, 255, 255)
    C_WHITE  = (30,  41,  59)
    C_BULLET = (71,  85,  105)
    C_ACCENT = (99,  102, 241)
    C_DIM    = (148, 163, 184)

    def _rgb(self, t: tuple) -> "_RGBColor":  # type: ignore[name-defined]
        return _RGBColor(t[0], t[1], t[2])

    def _add_rect(self, slide, left, top, width, height, color: tuple) -> None:
        try:
            shape = slide.shapes.add_shape(1, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(color)
            shape.line.fill.background()
        except Exception as exc:
            logger.debug("add_shape: %s", exc)

    def _add_text(
        self, slide, left, top, width, height,
        text: str, size: int, bold: bool, color: tuple, wrap: bool = True,
        font_name: str = "Arial"
    ) -> None:
        try:
            box = slide.shapes.add_textbox(left, top, width, height)
            tf  = box.text_frame
            tf.word_wrap = wrap
            para = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            run  = para.add_run()
            run.text           = str(text or "")
            run.font.name      = font_name
            run.font.size      = _Pt(size)
            run.font.bold      = bold
            run.font.color.rgb = self._rgb(color)
        except Exception as exc:
            logger.debug("add_text: %s", exc)

    def _add_bullets(self, slide, left, top, width, height, bullets: List[str]) -> None:
        try:
            box = slide.shapes.add_textbox(left, top, width, height)
            tf  = box.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                para = tf.paragraphs[0] if (i == 0 and tf.paragraphs) else tf.add_paragraph()
                para.space_before = _Pt(8)
                run  = para.add_run()
                run.text           = f"\u25b8  {str(bullet or '')}"
                run.font.name      = "Calibri"
                run.font.size      = _Pt(24) # Content font size (18-24pt)
                run.font.bold      = False
                run.font.color.rgb = self._rgb(self.C_BULLET)
        except Exception as exc:
            logger.debug("add_bullets: %s", exc)

    def _add_arrow(self, slide, sx, sy, ex, ey) -> None:
        try:
            # 1 is MSO_CONNECTOR.STRAIGHT
            conn = slide.shapes.add_connector(1, _Inches(sx), _Inches(sy), _Inches(ex), _Inches(ey))
            conn.line.color.rgb = self._rgb(self.C_ACCENT)
            conn.line.width = _Pt(1.5)
            try:
                line_elem = conn.line._get_or_add_ln()
                line_elem.append(_parse_xml(
                    '<a:tailEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
                ))
            except Exception as e:
                logger.debug("arrowhead xml error: %s", e)
        except Exception as exc:
            logger.debug("add_connector: %s", exc)

    def generate(self, items: list, deck_title: str = "Bài giảng", template_path: str | None = None) -> bytes:
        if not PPTX_OK:
            raise RuntimeError(f"python-pptx init failed: {PPTX_ERROR}")

        # Always enforce the premium built-in layouts to guarantee design consistency and stable diagram rendering
        template_path = None

        if template_path:
            prs = _Presentation(template_path)
            # Try to use layout 1 (usually Title and Content), fallback to 0 or 6
            layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
            slide_layout = prs.slide_layouts[layout_idx]

            for idx, item in enumerate(items):
                try:
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Fill Title
                    if slide.shapes.title:
                        slide.shapes.title.text = item.title
                    
                    # Fill Body Placeholders
                    body_shape = None
                    for shape in slide.placeholders:
                        if shape != slide.shapes.title and shape.has_text_frame:
                            body_shape = shape
                            break
                            
                    if body_shape:
                        tf = body_shape.text_frame
                        tf.text = "" # Clear default
                        for i, bullet in enumerate(item.bullet_points):
                            p = tf.paragraphs[0] if (i == 0 and tf.paragraphs) else tf.add_paragraph()
                            p.text = str(bullet)
                            # Let the template handle font and bullet styles automatically
                    
                    if item.speaker_notes and slide.has_notes_slide:
                        try:
                            slide.notes_slide.notes_text_frame.text = str(item.speaker_notes)
                        except Exception:
                            pass
                except Exception as exc:
                    logger.error("Slide %d build error with template: %s", idx, exc)

        else:
            # Default manual drawing
            prs = _Presentation()
            prs.slide_width  = _Inches(13.33)
            prs.slide_height = _Inches(7.5)
            blank = prs.slide_layouts[6]
            total = len(items)

            for idx, item in enumerate(items):
                try:
                    slide = prs.slides.add_slide(blank)
                    try:
                        bg = slide.background.fill
                        bg.solid()
                        bg.fore_color.rgb = self._rgb(self.C_BG)
                    except Exception:
                        pass

                    self._add_rect(slide, _Inches(0), _Inches(0), _Inches(0.12), _Inches(7.5), self.C_ACCENT)
                    self._add_text(slide, _Inches(0.3), _Inches(0.2), _Inches(12.7), _Inches(1.3), item.title, 40, True, self.C_WHITE, font_name="Arial")
                    self._add_rect(slide, _Inches(0.3), _Inches(1.6), _Inches(12.7), _Inches(0.045), self.C_ACCENT)

                    if item.diagram and item.diagram.nodes:
                        # Render two columns: left for bullets (first half), right for diagram
                        half = (len(item.bullet_points) + 1) // 2
                        left_bullets = item.bullet_points[:half]
                        
                        # Left bullets
                        self._add_bullets(slide, _Inches(0.45), _Inches(1.75), _Inches(5.8), _Inches(5.3), left_bullets)
                        
                        # Divider
                        self._add_rect(slide, _Inches(6.55), _Inches(1.6), _Inches(0.02), _Inches(5.3), self.C_DIM)
                        
                        # Right diagram (width=6.0, height=4.0, left=6.8, top=2.2)
                        nodes = item.diagram.nodes or []
                        links = item.diagram.links or []
                        N = len(nodes)
                        cx, cy = 180, 120
                        
                        # Calculate default positions if needed
                        base_positions = {}
                        if N == 1:
                            base_positions[nodes[0].id] = {"x": cx, "y": cy}
                        elif N == 2:
                            base_positions[nodes[0].id] = {"x": cx - 80, "y": cy}
                            base_positions[nodes[1].id] = {"x": cx + 80, "y": cy}
                        elif N == 3:
                            is_chain = len(links) == 2 and (
                                (links[0].source == nodes[0].id and links[0].target == nodes[1].id and links[1].source == nodes[1].id and links[1].target == nodes[2].id) or
                                (links[0].source == nodes[0].id and links[0].target == nodes[1].id and links[1].source == nodes[2].id and links[1].target == nodes[1].id)
                            )
                            if is_chain:
                                base_positions[nodes[0].id] = {"x": cx - 100, "y": cy}
                                base_positions[nodes[1].id] = {"x": cx, "y": cy}
                                base_positions[nodes[2].id] = {"x": cx + 100, "y": cy}
                            else:
                                rx = 90
                                ry = 60
                                for i, node in enumerate(nodes):
                                    angle = (i * 2 * math.pi) / N - math.pi / 2
                                    base_positions[node.id] = {
                                        "x": cx + rx * math.cos(angle),
                                        "y": cy + ry * math.sin(angle)
                                    }
                        else:
                            rx = 100
                            ry = 65
                            for i, node in enumerate(nodes):
                                angle = (i * 2 * math.pi) / N - math.pi / 2
                                base_positions[node.id] = {
                                    "x": cx + rx * math.cos(angle),
                                    "y": cy + ry * math.sin(angle)
                                }
                                
                        # Positions map (combine saved with base positions)
                        positions = {}
                        for node in nodes:
                            n_id = getattr(node, "id", None) or node.get("id") if isinstance(node, dict) else getattr(node, "id")
                            n_label = getattr(node, "label", None) or node.get("label") if isinstance(node, dict) else getattr(node, "label")
                            n_x = getattr(node, "x", None) or node.get("x") if isinstance(node, dict) else getattr(node, "x", None)
                            n_y = getattr(node, "y", None) or node.get("y") if isinstance(node, dict) else getattr(node, "y", None)
                            
                            if n_x is not None and n_y is not None:
                                positions[n_id] = {"x": n_x, "y": n_y, "label": n_label}
                            else:
                                base_pos = base_positions.get(n_id, {"x": cx, "y": cy})
                                positions[n_id] = {"x": base_pos["x"], "y": base_pos["y"], "label": n_label}

                        nodeW = 105
                        nodeH = 38

                        # Draw links (arrows) first so they are under the nodes
                        for link in links:
                            l_source = getattr(link, "source", None) or link.get("source") if isinstance(link, dict) else getattr(link, "source")
                            l_target = getattr(link, "target", None) or link.get("target") if isinstance(link, dict) else getattr(link, "target")
                            l_label = getattr(link, "label", "") or link.get("label", "") if isinstance(link, dict) else getattr(link, "label", "")
                            
                            from_pos = positions.get(l_source)
                            to_pos = positions.get(l_target)
                            if not from_pos or not to_pos:
                                continue

                            from_x, from_y = from_pos["x"], from_pos["y"]
                            to_x, to_y = to_pos["x"], to_pos["y"]

                            dx = to_x - from_x
                            dy = to_y - from_y
                            len_px = math.sqrt(dx * dx + dy * dy)
                            if len_px < 30:
                                continue

                            ux = dx / len_px
                            uy = dy / len_px

                            # Boundary-based padding
                            w = nodeW / 2
                            h = nodeH / 2
                            margin = 6

                            px_pad = w / abs(ux) if ux != 0 else float('inf')
                            py_pad = h / abs(uy) if uy != 0 else float('inf')
                            start_pad = min(px_pad, py_pad) + margin
                            end_pad = min(px_pad, py_pad) + 12 # 6px margin + 6px arrowhead

                            if len_px < start_pad + end_pad + 2:
                                continue

                            sx = from_x + ux * start_pad
                            sy = from_y + uy * start_pad
                            ex = to_x - ux * end_pad
                            ey = to_y - uy * end_pad

                            # Scale to Slide Inches (Right Column: left=6.8, top=2.2, width=6.0, height=4.0)
                            sx_in = 6.8 + sx / 60.0
                            sy_in = 2.2 + sy / 60.0
                            ex_in = 6.8 + ex / 60.0
                            ey_in = 2.2 + ey / 60.0

                            # Draw straight connector shape (type 1)
                            self._add_arrow(slide, sx_in, sy_in, ex_in, ey_in)

                            # Draw link label if any
                            if l_label:
                                nx = -uy
                                ny = ux
                                line_len = len_px - start_pad - end_pad
                                mx = sx + ux * line_len * 0.35 + nx * 10
                                my = sy + uy * line_len * 0.35 + ny * 10

                                mx_in = 6.8 + mx / 60.0
                                my_in = 2.2 + my / 60.0

                                label_text = str(l_label)
                                label_w_in = max(0.8, len(label_text) * 0.08 + 0.15)
                                label_h_in = 0.22

                                # Draw white background box with light border
                                try:
                                    rect = slide.shapes.add_shape(1, _Inches(mx_in - label_w_in/2), _Inches(my_in - label_h_in/2), _Inches(label_w_in), _Inches(label_h_in))
                                    rect.fill.solid()
                                    rect.fill.fore_color.rgb = self._rgb((255, 255, 255))
                                    rect.line.color.rgb = self._rgb((226, 232, 240)) # #e2e8f0
                                    rect.line.width = _Pt(0.5)
                                except Exception:
                                    pass

                                # Draw label text
                                self._add_text(slide, _Inches(mx_in - label_w_in/2), _Inches(my_in - label_h_in/2 - 0.02), _Inches(label_w_in), _Inches(label_h_in), label_text, 7, True, (79, 70, 229), wrap=False)

                        # Render nodes on top
                        for n_id, pos in positions.items():
                            nx_in = 6.8 + (pos["x"] - nodeW / 2) / 60.0
                            ny_in = 2.2 + (pos["y"] - nodeH / 2) / 60.0
                            w_in = nodeW / 60.0
                            h_in = nodeH / 60.0

                            # Draw node shape (rounded rectangle = type 5)
                            try:
                                rect = slide.shapes.add_shape(5, _Inches(nx_in), _Inches(ny_in), _Inches(w_in), _Inches(h_in))
                                rect.fill.solid()
                                rect.fill.fore_color.rgb = self._rgb((255, 255, 255))
                                rect.line.color.rgb = self._rgb(self.C_ACCENT)
                                rect.line.width = _Pt(1.5)
                                tf = rect.text_frame
                                tf.word_wrap = True
                                tf.margin_left = _Pt(2)
                                tf.margin_right = _Pt(2)
                                tf.margin_top = _Pt(2)
                                tf.margin_bottom = _Pt(2)
                                para = tf.paragraphs[0]
                                para.alignment = 2 # Center alignment (PP_ALIGN.CENTER)
                                run = para.add_run()
                                run.text = str(pos["label"])
                                run.font.name = "Arial"
                                run.font.size = _Pt(9)
                                run.font.bold = True
                                run.font.color.rgb = self._rgb((30, 41, 59)) # #1e293b
                            except Exception as node_exc:
                                logger.error("Node draw error: %s", node_exc)

                    else:
                        # Standard full-width layout
                        self._add_bullets(slide, _Inches(0.45), _Inches(1.75), _Inches(12.4), _Inches(5.3), item.bullet_points)

                    self._add_text(slide, _Inches(11.5), _Inches(7.1), _Inches(1.7), _Inches(0.35), f"{idx+1}/{total}", 11, False, self.C_DIM, wrap=False)

                    if item.speaker_notes:
                        try:
                            slide.notes_slide.notes_text_frame.text = str(item.speaker_notes)
                        except Exception:
                            pass
                except Exception as exc:
                    logger.error("Slide %d build error: %s", idx, exc)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()


# ── PDF Renderer ──────────────────────────────────────────────────────────────

class PdfRenderer:
    """Generates .pdf bytes (landscape A4) from a list of SlideItem objects."""

    PAGE_W, PAGE_H = (landscape(A4) if PDF_OK else (842, 595))
    BG     = HexColor("#ffffff") if PDF_OK else None
    ACCENT = HexColor("#6366f1") if PDF_OK else None
    WHITE  = HexColor("#1e293b") if PDF_OK else None
    BULLET = HexColor("#475569") if PDF_OK else None
    DIM    = HexColor("#94a3b8") if PDF_OK else None

    def _type_color(self, idx: int, total: int):
        if idx == 0:              return HexColor("#818cf8")
        if idx == total - 1 > 1: return HexColor("#4ade80")
        return HexColor("#38bdf8")

    def _type_label(self, idx: int, total: int) -> str:
        if idx == 0:              return "TITLE"
        if idx == total - 1 > 1: return "SUMMARY"
        return "CONTENT"

    def generate(self, items: list, deck_title: str = "Bài giảng") -> bytes:
        if not PDF_OK:
            raise RuntimeError(f"reportlab not available: {PDF_ERROR}")

        buf   = io.BytesIO()
        W, H  = self.PAGE_W, self.PAGE_H
        c     = _Canvas(buf, pagesize=(W, H))
        total = len(items)

        for idx, item in enumerate(items):
            try:
                c.setFillColor(self.BG)
                c.rect(0, 0, W, H, fill=1, stroke=0)

                c.setFillColor(self.ACCENT)
                c.rect(0, 0, 10, H, fill=1, stroke=0)

                c.setFillColor(self._type_color(idx, total))
                c.setFont("Helvetica-Bold", 8)
                c.drawRightString(W - 16, H - 20, f"{self._type_label(idx, total)}  {idx+1}/{total}")

                c.setFillColor(self.WHITE)
                c.setFont("Helvetica-Bold", 32)
                c.drawString(26, H - 68, str(item.title or "")[:100])

                c.setStrokeColor(self.ACCENT)
                c.setLineWidth(2.5)
                c.line(26, H - 82, W - 26, H - 82)

                if item.diagram and item.diagram.nodes:
                    # Render two columns: left for bullets (first half), right for diagram
                    half = (len(item.bullet_points) + 1) // 2
                    left_bullets = item.bullet_points[:half]
                    
                    c.setFont("Helvetica", 18)
                    y = H - 106
                    for bullet in left_bullets[:6]:
                        if y < 60:
                            break
                        c.setFillColor(self.ACCENT)
                        c.drawString(26, y, "▸")
                        c.setFillColor(self.BULLET)
                        c.drawString(42, y, str(bullet or "")[:120])
                        y -= 34
                        
                    # Divider in the middle
                    c.setStrokeColor(self.DIM)
                    c.setLineWidth(1.0)
                    c.line(421, H - 82, 421, 60)
                    
                    # Right column diagram area: width=360, height=240, left=446, top=135 in PDF
                    # scale = 1.0 (1-to-1 pixels)
                    # SVG y=0 is at PDF y=375, SVG y=240 is at PDF y=135
                    nodes = item.diagram.nodes or []
                    links = item.diagram.links or []
                    N = len(nodes)
                    cx, cy = 180, 120
                    
                    # Calculate default circular/linear positions if needed
                    base_positions = {}
                    if N == 1:
                        base_positions[nodes[0].id] = {"x": cx, "y": cy}
                    elif N == 2:
                        base_positions[nodes[0].id] = {"x": cx - 80, "y": cy}
                        base_positions[nodes[1].id] = {"x": cx + 80, "y": cy}
                    elif N == 3:
                        is_chain = len(links) == 2 and (
                            (links[0].source == nodes[0].id and links[0].target == nodes[1].id and links[1].source == nodes[1].id and links[1].target == nodes[2].id) or
                            (links[0].source == nodes[0].id and links[0].target == nodes[1].id and links[1].source == nodes[2].id and links[1].target == nodes[1].id)
                        )
                        if is_chain:
                            base_positions[nodes[0].id] = {"x": cx - 100, "y": cy}
                            base_positions[nodes[1].id] = {"x": cx, "y": cy}
                            base_positions[nodes[2].id] = {"x": cx + 100, "y": cy}
                        else:
                            rx = 90
                            ry = 60
                            for i, node in enumerate(nodes):
                                angle = (i * 2 * math.pi) / N - math.pi / 2
                                base_positions[node.id] = {
                                    "x": cx + rx * math.cos(angle),
                                    "y": cy + ry * math.sin(angle)
                                }
                    else:
                        rx = 100
                        ry = 65
                        for i, node in enumerate(nodes):
                            angle = (i * 2 * math.pi) / N - math.pi / 2
                            base_positions[node.id] = {
                                "x": cx + rx * math.cos(angle),
                                "y": cy + ry * math.sin(angle)
                            }
                            
                    # Combine saved positions with base positions
                    positions = {}
                    for node in nodes:
                        n_id = getattr(node, "id", None) or node.get("id") if isinstance(node, dict) else getattr(node, "id")
                        n_label = getattr(node, "label", None) or node.get("label") if isinstance(node, dict) else getattr(node, "label")
                        n_x = getattr(node, "x", None) or node.get("x") if isinstance(node, dict) else getattr(node, "x", None)
                        n_y = getattr(node, "y", None) or node.get("y") if isinstance(node, dict) else getattr(node, "y", None)
                        
                        if n_x is not None and n_y is not None:
                            positions[n_id] = {"x": n_x, "y": n_y, "label": n_label}
                        else:
                            base_pos = base_positions.get(n_id, {"x": cx, "y": cy})
                            positions[n_id] = {"x": base_pos["x"], "y": base_pos["y"], "label": n_label}

                    nodeW = 105
                    nodeH = 38

                    # Draw links (arrows) first so they are under the nodes
                    for link in links:
                        l_source = getattr(link, "source", None) or link.get("source") if isinstance(link, dict) else getattr(link, "source")
                        l_target = getattr(link, "target", None) or link.get("target") if isinstance(link, dict) else getattr(link, "target")
                        l_label = getattr(link, "label", "") or link.get("label", "") if isinstance(link, dict) else getattr(link, "label", "")
                        
                        from_pos = positions.get(l_source)
                        to_pos = positions.get(l_target)
                        if not from_pos or not to_pos:
                            continue

                        from_x, from_y = from_pos["x"], from_pos["y"]
                        to_x, to_y = to_pos["x"], to_pos["y"]

                        dx = to_x - from_x
                        dy = to_y - from_y
                        len_px = math.sqrt(dx * dx + dy * dy)
                        if len_px < 30:
                            continue

                        ux = dx / len_px
                        uy = dy / len_px

                        # Boundary-based padding
                        w = nodeW / 2
                        h = nodeH / 2
                        margin = 6

                        px_pad = w / abs(ux) if ux != 0 else float('inf')
                        py_pad = h / abs(uy) if uy != 0 else float('inf')
                        start_pad = min(px_pad, py_pad) + margin
                        end_pad = min(px_pad, py_pad) + 12 # 6px margin + 6px arrowhead

                        if len_px < start_pad + end_pad + 2:
                            continue

                        sx = from_x + ux * start_pad
                        sy = from_y + uy * start_pad
                        ex = to_x - ux * end_pad
                        ey = to_y - uy * end_pad

                        # Scale to PDF page coordinates: x_pdf = 446 + x, y_pdf = 375 - y
                        sx_pdf = 446 + sx
                        sy_pdf = 375 - sy
                        ex_pdf = 446 + ex
                        ey_pdf = 375 - ey

                        # Draw link line
                        c.setStrokeColor(self.ACCENT)
                        c.setLineWidth(1.5)
                        c.line(sx_pdf, sy_pdf, ex_pdf, ey_pdf)

                        # Draw arrowhead (triangle pointing to ex_pdf, ey_pdf)
                        # We calculate directional vectors in PDF space (origin bottom-left)
                        dx_pdf = ex_pdf - sx_pdf
                        dy_pdf = ey_pdf - sy_pdf
                        len_pdf = math.sqrt(dx_pdf * dx_pdf + dy_pdf * dy_pdf)
                        if len_pdf > 0:
                            ux_p = dx_pdf / len_pdf
                            uy_p = dy_pdf / len_pdf
                            nx_p = -uy_p
                            ny_p = ux_p

                            cx1 = ex_pdf - ux_p * 8 + nx_p * 3
                            cy1 = ey_pdf - uy_p * 8 + ny_p * 3
                            cx2 = ex_pdf - ux_p * 8 - nx_p * 3
                            cy2 = ey_pdf - uy_p * 8 - ny_p * 3

                            arrow_path = c.beginPath()
                            arrow_path.moveTo(ex_pdf, ey_pdf)
                            arrow_path.lineTo(cx1, cy1)
                            arrow_path.lineTo(cx2, cy2)
                            arrow_path.close()
                            c.setFillColor(self.ACCENT)
                            c.drawPath(arrow_path, fill=1, stroke=0)

                        # Draw link label if any
                        if l_label:
                            nx = -uy
                            ny = ux
                            line_len = len_px - start_pad - end_pad
                            mx = sx + ux * line_len * 0.35 + nx * 10
                            my = sy + uy * line_len * 0.35 + ny * 10

                            mx_pdf = 446 + mx
                            my_pdf = 375 - my

                            label_text = str(l_label)
                            label_w = max(40, len(label_text) * 5 + 8)

                            # White background box with light border
                            c.setFillColor(HexColor("#ffffff"))
                            c.setStrokeColor(HexColor("#e2e8f0"))
                            c.setLineWidth(0.5)
                            c.rect(mx_pdf - label_w/2, my_pdf - 6, label_w, 12, fill=1, stroke=1)

                            # Label text
                            c.setFillColor(HexColor("#4f46e5"))
                            c.setFont("Helvetica-Bold", 7)
                            c.drawCentredString(mx_pdf, my_pdf - 2, label_text)

                    # Render nodes on top
                    for n_id, pos in positions.items():
                        nx_pdf = 446 + pos["x"]
                        ny_pdf = 375 - pos["y"]

                        # Draw node (rounded rectangle)
                        c.setFillColor(HexColor("#ffffff"))
                        c.setStrokeColor(self.ACCENT)
                        c.setLineWidth(1.5)
                        c.roundRect(nx_pdf - nodeW/2, ny_pdf - nodeH/2, nodeW, nodeH, 6, fill=1, stroke=1)

                        # Draw node label
                        c.setFillColor(HexColor("#1e293b"))
                        c.setFont("Helvetica-Bold", 8)
                        c.drawCentredString(nx_pdf, ny_pdf - 3, str(pos["label"]))

                else:
                    # Standard full-width layout
                    c.setFont("Helvetica", 18)
                    y = H - 106
                    for bullet in item.bullet_points[:6]:
                        if y < 60:
                            break
                        c.setFillColor(self.ACCENT)
                        c.drawString(26, y, "▸")
                        c.setFillColor(self.BULLET)
                        c.drawString(42, y, str(bullet or "")[:120])
                        y -= 34

                c.setFillColor(self.DIM)
                c.setFont("Helvetica", 9)
                c.drawCentredString(W / 2, 16, deck_title[:60])

            except Exception as ex:
                logger.debug("PDF slide %d: %s", idx, ex)

            c.showPage()

        c.save()
        return buf.getvalue()


# ── Singleton instances (module-level) ────────────────────────────────────────

pptx_renderer = PptxRenderer()
pdf_renderer  = PdfRenderer()
